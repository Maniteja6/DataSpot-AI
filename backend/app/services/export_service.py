"""
Generates real, downloadable export artifacts from a dataset's cleaned
DataFrame plus its computed insights/decisions: CSV, JSON, and Excel via
pandas; a PDF executive report via reportlab; a PowerPoint summary deck via
python-pptx.
"""

from __future__ import annotations
import json
import uuid
from pathlib import Path
from app.config.settings import get_settings
from app.models.dataset import Dataset
from app.models.insight import Insight
from app.models.decision import DecisionCard
from app.services.dataset_service import get_cleaned_dataframe

EXPORT_DIR = Path("./storage/exports")
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


def _export_path(dataset_id: str, fmt: str) -> Path:
    ext = {"excel": "xlsx"}.get(fmt, fmt)
    return EXPORT_DIR / f"{dataset_id}_{uuid.uuid4().hex[:8]}.{ext}"


def export_csv(dataset: Dataset) -> Path:
    df = get_cleaned_dataframe(dataset.id)
    path = _export_path(dataset.id, "csv")
    (df if df is not None else _empty_frame()).to_csv(path, index=False)
    return path


def export_json(dataset: Dataset, insights: list[Insight], decisions: list[DecisionCard]) -> Path:
    df = get_cleaned_dataframe(dataset.id)
    path = _export_path(dataset.id, "json")
    payload = {
        "dataset": json.loads(dataset.model_dump_json()),
        "insights": [json.loads(i.model_dump_json()) for i in insights],
        "decisions": [json.loads(d.model_dump_json()) for d in decisions],
        "rows_preview": df.head(50).to_dict(orient="records") if df is not None else [],
    }
    path.write_text(json.dumps(payload, indent=2))
    return path


def export_excel(dataset: Dataset, insights: list[Insight]) -> Path:
    import pandas as pd

    df = get_cleaned_dataframe(dataset.id)
    if df is None:
        df = _empty_frame()
    path = _export_path(dataset.id, "excel")
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        df.to_excel(writer, sheet_name="Cleaned Data", index=False)
        insights_df = pd.DataFrame([i.model_dump() for i in insights])
        (insights_df if not insights_df.empty else pd.DataFrame([{"note": "No insights generated yet"}])).to_excel(
            writer, sheet_name="Insights", index=False
        )
    return path


def export_pdf(dataset: Dataset, insights: list[Insight], decisions: list[DecisionCard], summary_text: str) -> Path:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    path = _export_path(dataset.id, "pdf")
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=LETTER)
    story = [
        Paragraph(f"DataSpot AI — Executive Report: {dataset.name}", styles["Title"]),
        Spacer(1, 0.2 * inch),
        Paragraph(summary_text or "No executive summary generated yet.", styles["BodyText"]),
        Spacer(1, 0.25 * inch),
        Paragraph("Key Findings", styles["Heading2"]),
    ]
    for insight in insights[:8]:
        story.append(Paragraph(f"• {insight.title}: {insight.narrative}", styles["BodyText"]))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Strategic Recommendations", styles["Heading2"]))
    for decision in decisions[:8]:
        story.append(
            Paragraph(
                f"• [{decision.priority.value.upper()}] {decision.title} "
                f"(expected ROI {decision.expected_roi_pct}%)",
                styles["BodyText"],
            )
        )
    doc.build(story)
    return path


def export_pptx(dataset: Dataset, insights: list[Insight], decisions: list[DecisionCard]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    path = _export_path(dataset.id, "pptx")
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"DataSpot AI — {dataset.name}"
    title_slide.placeholders[1].text = "Executive Summary Deck"

    bullet_layout = prs.slide_layouts[1]

    findings_slide = prs.slides.add_slide(bullet_layout)
    findings_slide.shapes.title.text = "Key Findings"
    body = findings_slide.placeholders[1].text_frame
    body.clear()
    for i, insight in enumerate(insights[:6]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = f"{insight.title}"
        p.font.size = Pt(18)

    decisions_slide = prs.slides.add_slide(bullet_layout)
    decisions_slide.shapes.title.text = "Strategic Recommendations"
    body = decisions_slide.placeholders[1].text_frame
    body.clear()
    for i, decision in enumerate(decisions[:6]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = f"[{decision.priority.value.upper()}] {decision.title}"
        p.font.size = Pt(18)

    prs.save(str(path))
    return path


def _empty_frame():
    import pandas as pd

    return pd.DataFrame()
